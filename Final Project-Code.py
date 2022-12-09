#!/usr/bin/env python
# coding: utf-8

# In[1]:


import os
import openai
import requests
import json

from pptx import Presentation 
get_ipython().run_line_magic('env', 'OPENAI_API_KEY=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')


openai.api_key = os.getenv("OPENAI_API_KEY")

# from huggingface_hub import notebook_login
# notebook_login()

# import torch
# from torch import autocast
# from diffusers import StableDiffusionPipeline
# from PIL import Image


# In[2]:


# Generate structure

def structure_(name):
    
    
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt="Create a story structure for"+ name,
        temperature=0.7,
        max_tokens=100,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
        )
## Title: Essay outline
## Availibility: https://beta.openai.com/examples/default-essay-outline

    struc_text = response.choices[0].text

    for i in range(len(struc_text)):
                   
        sub_topic_list = struc_text.split("\n\n")
        
        sub_topic_title = []
        sub_topic_des = []
        
        for i in sub_topic_list:
            
            if i == '':
                
                sub_topic_list.remove(i)

        
        for t in sub_topic_list:                
            
            sub_topic_split = t.split(":")
#             print(sub_topic_split)
            sub_topic_title.append(sub_topic_split[0])
            sub_topic_des.append(sub_topic_split[1])
#           ** add function to remove empty string
        topic = str(name)
        name = {}
        name['topic'] = topic
        name['subtopics'] = sub_topic_list
        name['title'] = sub_topic_title 
        name['des'] = sub_topic_des
        
        return name


# In[3]:


def gen_des(des):
    
        
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt= des + 'in bullet points',
        temperature=0.7,
        max_tokens=100,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )
    descriptions = response.choices[0].text

    des_text = descriptions[2:]
    
    return des_text


# In[4]:


def sum_text(name):
    
    name['des_text'] = des
    
    sum_text = []
    
    for d in des:
        
        desc_ = ''.join(i)
    
        response = openai.Completion.create(
            model="text-davinci-003",
            prompt="Summarize the following text in 10 words:"+ desc_,
            temperature=0.7,
            max_tokens=100,
            top_p=1,
            frequency_penalty=0,
            presence_penalty=0
            )
# From https://beta.openai.com/docs/api-reference/introduction

        sum_ = response.choices[0].text    
        sum_text.append(sum_)
        
    name['sum_text'] = sum_text
    
    return sum_text    


# In[ ]:


# from diffusers import StableDiffusionPipeline

# def gen_img(sum_text):
#     pipe = StableDiffusionPipeline.from_pretrained("runwayml/stable-diffusion-v1-5")
#     pipe = pipe.to("cpu")

#     # Recommended if your computer has < 64 GB of RAM
#     pipe.enable_attention_slicing()

#     prompt = sum_text

#     # First-time "warmup" pass (see explanation above)
#     _ = pipe(prompt, num_inference_steps=1)

#     # Results match those from the CPU device after the warmup pass.
#     image = pipe(prompt).images[0]

## *    Title: How to use Stable Diffusion in Apple Silicon (M1/M2)
## *    Availability: https://huggingface.co/docs/diffusers/optimization/mps

# @InProceedings{Rombach_2022_CVPR,
#     author    = {Rombach, Robin and Blattmann, Andreas and Lorenz, Dominik and Esser, Patrick and Ommer, Bj\"orn},
#     title     = {High-Resolution Image Synthesis With Latent Diffusion Models},
#     booktitle = {Proceedings of the IEEE/CVF Conference on Computer Vision and Pattern Recognition (CVPR)},
#     month     = {June},
#     year      = {2022},
#     pages     = {10684-10695}
# }


# In[5]:


def gen_pre(name):
    
    pre = Presentation()

    first_slide_layout = pre.slide_layouts[0]
    rest_slide_layout = pre.slide_layouts[1]

    slide_0 = pre.slides.add_slide(first_slide_layout)
    slide_0.shapes.title.text = name['topic']
    
    for s in range(len(name['title'])):
        
        slide = pre.slides.add_slide(rest_slide_layout)
        slide.shapes.title.text = name['title'][s]
        slide.placeholders[1].text = name['des_text'][s]
        
#       for image generation  
#         slide_img = pre.slides.add_slide(rest_slide_layout)
#         slide.placeholders[1].insert_picture = gen_img(name['sum_text'][s])
            
    pre.save(name['topic'])


# In[6]:


def pro_(names):
    
    
    for n in names:
        
        n = structure_(n)
        
        des_text = []
        
        for i in range(len(n['des'])):
            
            
            des_ = gen_des(n['des'][i])
            des_text.append(des_)
#             print(des_text)
        
        n['des_text'] = des_text
        
#         print(n['des_text'])
        
        gen_pre(n)


# In[10]:


pro_(['Nikolas Tesla'])

